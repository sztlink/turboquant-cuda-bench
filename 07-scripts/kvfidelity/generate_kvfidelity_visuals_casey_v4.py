#!/usr/bin/env python3
import json, textwrap, re
from collections import Counter
from pathlib import Path

import pandas as pd
import plotly.graph_objects as go
from PIL import Image, ImageDraw, ImageFont

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
except Exception:
    Presentation = None

SRC = Path('/home/aya/implante/tmp/kvfidelity-aime24-n30-traces-2026-05-13.jsonl')
OUT = Path('/home/aya/implante/tmp/kvfidelity-casey-v4-2026-05-14')
OUT.mkdir(parents=True, exist_ok=True)

RUN_ORDER = [
    'fullkv_2048','tri_b256_2048','cask_b256_2048',
    'fullkv_4096','tri_b256_4096','cask_b256_4096',
    'tri_b384_4096','cask_b384_4096','tri_b512_4096','cask_b512_4096'
]
IMPORTANT = [0,7,9,11,12,24,26]

LABELS = {
    'clean_correct': {'sym':'✓','color':'#007A4D','name':'closed clean'},
    'correct_unboxed': {'sym':'u','color':'#69B99A','name':'closed weak'},
    'answer_marker_drift': {'sym':'D','color':'#D65A3B','name':'answer drift'},
    'latent_answer_marker_not_closed': {'sym':'A','color':'#C98713','name':'latent after marker'},
    'latent_final_zone_not_closed': {'sym':'F','color':'#E3B341','name':'latent final zone'},
    'latent_intermediate_only': {'sym':'~','color':'#7655B7','name':'intermediate only'},
    'evaluator_missed_boxed_gt': {'sym':'B','color':'#368CBF','name':'boxed missed'},
    'no_gt_seen': {'sym':'·','color':'#D6D1C7','name':'not discovered'},
}
LABEL_ORDER = list(LABELS.keys())
LABEL_VALUE = {k:i for i,k in enumerate(LABEL_ORDER)}
COLORSCALE=[]
for k,i in LABEL_VALUE.items():
    lo=i/(len(LABEL_ORDER)-1)
    COLORSCALE += [[lo,LABELS[k]['color']],[lo,LABELS[k]['color']]]

BG = '#FFFDF7'
PAPER = '#F7F3EA'
FG = '#151515'
MUTED = '#6F6A60'
GRID = '#D8D0C2'
ACCENT = '#007A4D'
ORANGE = '#D65A3B'
PURPLE = '#7655B7'
BLUE = '#368CBF'
YELLOW = '#C98713'

CLASS_COLORS = {
    'not_discovered':'#D6D1C7',
    'discovered_not_retained':PURPLE,
    'closure_failure':ORANGE,
    'closed':ACCENT,
    'extractor_bug':BLUE,
}


def load():
    rows=[json.loads(line) for line in SRC.read_text().splitlines() if line.strip()]
    df=pd.DataFrame(rows)
    df['run_id'] = pd.Categorical(df['run_id'], categories=RUN_ORDER, ordered=True)
    df['closed'] = df['eval_correct'].astype(bool)
    df['discovered'] = df['gt_anywhere'].astype(bool)
    df['retained'] = (df['gt_in_final_zone'].astype(bool) | df['gt_after_answer_marker'].astype(bool) | df['gt_boxed'].astype(bool))
    df['run_short'] = df['run_id'].astype(str).str.replace('_','<br>')
    return df.sort_values(['idx','run_id'])


def style(fig, title=None, w=1600, h=900):
    fig.update_layout(
        template='plotly_white', paper_bgcolor=PAPER, plot_bgcolor=BG,
        font=dict(family='Inter, Arial, sans-serif', color=FG, size=18),
        title=dict(text=title or '', x=0.02, y=0.98, xanchor='left', yanchor='top', font=dict(size=34, color=FG)),
        margin=dict(l=80,r=60,t=105,b=90), width=w, height=h,
    )
    fig.update_xaxes(gridcolor=GRID, zerolinecolor=GRID, tickfont=dict(color=MUTED, size=14))
    fig.update_yaxes(gridcolor=GRID, zerolinecolor=GRID, tickfont=dict(color=MUTED, size=14))
    return fig


def save(fig, name, w=1600, h=900):
    html=OUT/f'{name}.html'
    png=OUT/f'{name}.png'
    fig.write_html(str(html), include_plotlyjs='cdn')
    fig.write_image(str(png), width=w, height=h, scale=2)
    return png


def run_label(run):
    return run.replace('_','<br>')


def topology(df):
    mat=[]; text=[]; hover=[]
    for idx in range(30):
        row=[]; trow=[]; hrow=[]
        for run in RUN_ORDER:
            r=df[(df.idx==idx)&(df.run_id.astype(str)==run)].iloc[0]
            lab=r.pathology_label
            row.append(LABEL_VALUE[lab]); trow.append(LABELS[lab]['sym'])
            hrow.append(f"idx {idx}<br>GT {r.gt_norm}<br>{run}<br>{LABELS[lab]['name']}<br>pred {r.final_pred_norm}<br>first {r.first_gt_frac}<br>last {r.last_gt_frac}<br>flips {r.candidate_flips}")
        mat.append(row); text.append(trow); hover.append(hrow)
    fig=go.Figure(go.Heatmap(
        z=mat, x=[run_label(r) for r in RUN_ORDER], y=list(range(30)),
        colorscale=COLORSCALE, showscale=False, text=text, texttemplate='%{text}',
        textfont=dict(size=22, color='#151515'), hovertext=hover, hoverinfo='text', xgap=3, ygap=2
    ))
    fig=style(fig, 'Topology matrix: answer as temporal event', 1800, 1450)
    fig.update_layout(margin=dict(l=95,r=360,t=120,b=150))
    fig.update_yaxes(autorange='reversed', title='AIME24 problem index')
    fig.update_xaxes(title='runtime condition', tickfont=dict(size=13, color=MUTED), title_standoff=28)
    fig.add_annotation(text='Legend', x=1.025, y=0.98, xref='paper', yref='paper', showarrow=False, xanchor='left', font=dict(size=17,color=FG))
    for i,(k,v) in enumerate(LABELS.items()):
        fig.add_annotation(text=f"<span style='color:{v['color']}; font-size:22px'>{v['sym']}</span> {v['name']}",
            x=1.025, y=0.925 - i*0.055, xref='paper', yref='paper', showarrow=False, xanchor='left', align='left', font=dict(size=15,color=MUTED))
    return save(fig,'01_topology_matrix',1800,1450)


def drc_cycle(df):
    """Casey intervention: replace generic grouped bars with process strips Discovery → Retention → Closure."""
    fig=go.Figure()
    yvals=list(range(len(RUN_ORDER)))
    stage_x=[0,1,2]
    stage_names=['Discovery','Retention','Closure']
    stage_colors=[PURPLE,YELLOW,ACCENT]
    counts=[]
    for run in RUN_ORDER:
        sub=df[df.run_id.astype(str)==run]
        counts.append([
            int(sub.discovered.sum()),
            int(sub.retained.sum()),
            int(sub.closed.sum())
        ])
    # background guides
    for yi,run in enumerate(RUN_ORDER):
        fig.add_trace(go.Scatter(x=[0,2], y=[yi,yi], mode='lines', line=dict(color=GRID,width=3), showlegend=False, hoverinfo='skip'))
        # loss bands between stages: line width = loss count
        d,r,c=counts[yi]
        losses=[max(0,d-r), max(0,r-c)]
        for j,loss in enumerate(losses):
            if loss:
                fig.add_trace(go.Scatter(x=[j,j+1], y=[yi,yi], mode='lines', line=dict(color=hex_to_rgba(ORANGE,0.28), width=5+loss*2), showlegend=False, hoverinfo='skip'))
        for j,val in enumerate([d,r,c]):
            size=16+val*1.8
            fig.add_trace(go.Scatter(
                x=[j], y=[yi], mode='markers+text', text=[str(val)], textposition='middle center',
                marker=dict(size=size, color=stage_colors[j], line=dict(color=FG,width=1)),
                textfont=dict(size=14, color='#FFFDF7' if val>0 else FG),
                hovertext=f"{run}<br>{stage_names[j]}: {val}/30", hoverinfo='text', showlegend=False
            ))
    fig=style(fig, 'Answer lifecycle: where the answer is lost', 1700, 1050)
    fig.update_layout(margin=dict(l=260,r=80,t=125,b=155))
    fig.update_xaxes(tickmode='array', tickvals=stage_x, ticktext=stage_names, range=[-0.35,2.35], title='answer lifecycle stage', title_standoff=26)
    fig.update_yaxes(tickmode='array', tickvals=yvals, ticktext=[r.replace('_',' ') for r in RUN_ORDER], autorange='reversed', title='')
    # Keep explanatory text inside the safe plot area so it does not get cut in PPT/PDF export.
    fig.add_annotation(text='node label = count out of 30 · orange band = loss between stages', x=0.5, y=-0.105, xref='paper', yref='paper', showarrow=False, xanchor='center', font=dict(size=16,color=MUTED))
    return save(fig,'02_answer_lifecycle_strips',1700,1050)


def drc_summary_appendix(df):
    runs=[]; discovery=[]; retention=[]; closure=[]
    for run in RUN_ORDER:
        sub=df[df.run_id.astype(str)==run]
        runs.append(run_label(run)); discovery.append(sub.discovered.mean()*100); retention.append(sub.retained.mean()*100); closure.append(sub.closed.mean()*100)
    fig=go.Figure()
    fig.add_bar(name='Discovery: GT appears', x=runs, y=discovery, marker_color=PURPLE)
    fig.add_bar(name='Retention: GT near final/marker', x=runs, y=retention, marker_color=YELLOW)
    fig.add_bar(name='Closure: eval correct', x=runs, y=closure, marker_color=ACCENT)
    fig=style(fig,'Appendix: Discovery / Retention / Closure counts',1800,950)
    fig.update_layout(barmode='group', legend=dict(orientation='h', y=1.05, x=0.02), yaxis_ticksuffix='%')
    fig.update_yaxes(range=[0,45], title='share of 30 problems')
    return save(fig,'A01_drc_bars',1800,950)


def pathology_stacked(df):
    fig=go.Figure()
    for lab in LABEL_ORDER:
        ys=[]
        for run in RUN_ORDER:
            sub=df[df.run_id.astype(str)==run]
            ys.append((sub.pathology_label==lab).sum())
        fig.add_bar(name=LABELS[lab]['name'], x=[run_label(r) for r in RUN_ORDER], y=ys, marker_color=LABELS[lab]['color'])
    fig=style(fig,'Appendix: pathology mix',1800,950)
    fig.update_layout(barmode='stack', margin=dict(l=90,r=50,t=135,b=175), legend=dict(orientation='h', y=-0.22, x=0.0, xanchor='left', yanchor='top', font=dict(size=14), traceorder='normal', itemwidth=120))
    fig.update_yaxes(title='count of 30 samples')
    return save(fig,'A02_pathology_stack',1800,950)


def hex_to_rgba(hex_color, alpha=0.64):
    h=hex_color.lstrip('#')
    r=int(h[0:2],16); g=int(h[2:4],16); b=int(h[4:6],16)
    return f'rgba({r},{g},{b},{alpha})'


def sankey_for(df, a, b, name, title):
    labels=['not_discovered','discovered_not_retained','closure_failure','closed']
    left=[f'2048<br>{l}' for l in labels]
    right=[f'4096<br>{l}' for l in labels]
    node_labels=left+right
    node_colors=[CLASS_COLORS[l] for l in labels]+[CLASS_COLORS[l] for l in labels]
    trans=Counter()
    for idx in range(30):
        ra=df[(df.idx==idx)&(df.run_id.astype(str)==a)].iloc[0]
        rb=df[(df.idx==idx)&(df.run_id.astype(str)==b)].iloc[0]
        trans[(ra.pathology_class, rb.pathology_class)] += 1
    source=[]; target=[]; value=[]; link_color=[]
    for (x,y),v in trans.items():
        source.append(labels.index(x)); target.append(len(labels)+labels.index(y)); value.append(v); link_color.append(hex_to_rgba(CLASS_COLORS[y], 0.62))
    fig=go.Figure(go.Sankey(
        arrangement='fixed',
        node=dict(label=node_labels, color=node_colors, pad=18, thickness=24, line=dict(color='#FFFDF7',width=1)),
        link=dict(source=source,target=target,value=value,color=link_color)
    ))
    fig=style(fig, title, 1500, 900)
    return save(fig, name,1500,900)


def transitions(df):
    return [
        sankey_for(df,'fullkv_2048','fullkv_4096','04_sankey_fullkv_2048_to_4096','FullKV: answer states transform from 2048 → 4096'),
        sankey_for(df,'tri_b256_2048','tri_b256_4096','A03_sankey_tri_b256_2048_to_4096','TriAttention b256: 2048 → 4096 transition'),
        sankey_for(df,'cask_b256_2048','cask_b256_4096','A04_sankey_cask_b256_2048_to_4096','CASK b256: 2048 → 4096 transition'),
    ]


def scatter(df):
    fig=go.Figure()
    for lab in LABEL_ORDER:
        sub=df[df.pathology_label==lab]
        if sub.empty: continue
        fig.add_trace(go.Scatter(
            x=sub.output_chars, y=sub.candidate_flips, mode='markers', name=LABELS[lab]['name'],
            text=[f"{r.run_id}<br>idx {r.idx}<br>gt {r.gt_norm}<br>{LABELS[lab]['name']}" for r in sub.itertuples()],
            hoverinfo='text+x+y', marker=dict(size=12, color=LABELS[lab]['color'], line=dict(width=1,color='#000'))
        ))
    fig=style(fig,'Appendix: candidate churn vs output length',1600,950)
    fig.update_xaxes(title='output characters')
    fig.update_yaxes(title='candidate flips near answer markers')
    return save(fig,'A05_candidate_churn_scatter',1600,950)


def case_strip(df, idx):
    """Casey intervention: pred/flips printed, duration band between first/last GT, churn cue visible without hover."""
    sub=df[df.idx==idx].copy()
    fig=go.Figure()
    ylabels=[r.replace('_',' ') for r in RUN_ORDER]
    ymap={run:i for i,run in enumerate(RUN_ORDER)}
    max_flips=max(1, int(df.candidate_flips.max()))
    for run in RUN_ORDER:
        y=ymap[run]
        fig.add_trace(go.Scatter(x=[0,1],y=[y,y],mode='lines',line=dict(color=GRID,width=4),showlegend=False,hoverinfo='skip'))
    for _,r in sub.iterrows():
        y=ymap[str(r.run_id)]
        lab=r.pathology_label
        color=LABELS[lab]['color']
        first = None if pd.isna(r.first_gt_frac) else float(r.first_gt_frac)
        last = None if pd.isna(r.last_gt_frac) else float(r.last_gt_frac)
        flips=int(r.candidate_flips or 0)
        # candidate churn density, visible as small ticks near the right side
        tick_count=min(14, max(0, round(flips/max_flips*14)))
        for k in range(tick_count):
            x=0.78 + k*0.014
            fig.add_trace(go.Scatter(x=[x,x], y=[y-0.18,y-0.06], mode='lines', line=dict(color=hex_to_rgba(ORANGE,0.48), width=1.2), showlegend=False, hoverinfo='skip'))
        # duration between first and last occurrence
        if first is not None and last is not None and last >= first:
            fig.add_trace(go.Scatter(x=[first,last],y=[y,y],mode='lines',line=dict(color=hex_to_rgba(color,0.38),width=16),showlegend=False,hoverinfo='skip'))
        if first is not None:
            fig.add_trace(go.Scatter(x=[first],y=[y],mode='markers',marker=dict(symbol='triangle-up',size=17,color=color,line=dict(color=FG,width=1)),showlegend=False,hovertext=f"{r.run_id}<br>first GT {first:.3f}",hoverinfo='text'))
        if last is not None:
            fig.add_trace(go.Scatter(x=[last],y=[y],mode='markers',marker=dict(symbol='diamond',size=17,color=color,line=dict(color=FG,width=1)),showlegend=False,hovertext=f"{r.run_id}<br>last GT {last:.3f}",hoverinfo='text'))
        fig.add_trace(go.Scatter(x=[1.025],y=[y],mode='markers+text',text=[LABELS[lab]['sym']],textposition='middle center',marker=dict(symbol='square',size=28,color=color,line=dict(color=FG,width=1)),textfont=dict(size=18,color='#151515'),showlegend=False,hoverinfo='skip'))
        fig.add_annotation(text=f"pred {r.final_pred_norm} · flips {flips}", x=1.06, y=y, xref='x', yref='y', showarrow=False, xanchor='left', font=dict(size=12,color=MUTED))
    gt=sub.iloc[0].gt_norm
    fig=style(fig,f'Case strip idx {idx}: GT {gt}',1900,900)
    fig.update_layout(margin=dict(l=220,r=260,t=110,b=115))
    fig.update_xaxes(range=[-0.02,1.18], title='normalized position in output', tickformat='.0%')
    fig.update_yaxes(tickmode='array', tickvals=list(range(len(RUN_ORDER))), ticktext=ylabels, autorange='reversed')
    fig.add_annotation(text='▲ first GT · ◆ last GT · translucent band = GT persistence · orange microticks = candidate churn · square = final pathology', x=0.02, y=-0.16, xref='paper', yref='paper', showarrow=False, align='left', xanchor='left', font=dict(size=18,color=MUTED))
    return save(fig,f'case_strip_idx_{idx:02d}',1900,900)


def make_case_strips(df):
    return [case_strip(df,idx) for idx in IMPORTANT]


def font(path_candidates, size):
    for p in path_candidates:
        if Path(p).exists():
            return ImageFont.truetype(p, size)
    return ImageFont.load_default()


def draw_wrapped(draw, text, xy, font_obj, fill, width_chars, line_gap=6):
    x,y=xy
    lines=[]
    for para in text.split('\n'):
        if not para.strip():
            lines.append('')
        else:
            lines.extend(textwrap.wrap(para, width=width_chars, break_long_words=False, replace_whitespace=False))
    for line in lines:
        draw.text((x,y), line, font=font_obj, fill=fill)
        y += font_obj.size + line_gap
    return y


def label_birth_panel(df, idx=9, run='fullkv_4096'):
    """Casey intervention: one slide showing how a label is born from the trace/output."""
    r=df[(df.idx==idx)&(df.run_id.astype(str)==run)].iloc[0]
    W,H=1900,1080
    img=Image.new('RGB',(W,H),PAPER)
    d=ImageDraw.Draw(img)
    title_f=font(['/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf'],42)
    head_f=font(['/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf'],24)
    mono_f=font(['/usr/share/fonts/truetype/dejavu/DejaVuSansMono.ttf'],20)
    small_f=font(['/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf'],18)
    d.text((70,55),'How the label is born: idx 9 / FullKV 4096',font=title_f,fill=FG)
    d.text((72,112),'The atlas is strongest when the measurement process becomes visible, not only the result.',font=small_f,fill=MUTED)
    d.line((70,150,W-70,150),fill=GRID,width=2)
    # left metadata card
    card=(70,190,580,930)
    d.rounded_rectangle(card, radius=18, fill=BG, outline=GRID, width=2)
    d.text((105,225),'Trace signals',font=head_f,fill=FG)
    meta=[
        ('run', run), ('GT', str(r.gt_norm)), ('evaluator pred', str(r.final_pred_norm)),
        ('pathology', LABELS[r.pathology_label]['name']), ('first GT frac', str(round(float(r.first_gt_frac),3)) if pd.notna(r.first_gt_frac) else '—'),
        ('last GT frac', str(round(float(r.last_gt_frac),3)) if pd.notna(r.last_gt_frac) else '—'),
        ('candidate flips', str(int(r.candidate_flips))), ('closure', 'failed')
    ]
    y=275
    for k,v in meta:
        d.text((105,y),k,font=small_f,fill=MUTED)
        d.text((300,y),v,font=small_f,fill=FG)
        y+=48
    # mini grammar
    y+=32
    d.text((105,y),'Visual grammar',font=head_f,fill=FG); y+=45
    for sym,label,col in [('▲','first GT',PURPLE),('◆','last GT',PURPLE),('■','final label',LABELS[r.pathology_label]['color']),('▥','churn ticks',ORANGE)]:
        d.text((110,y),sym,font=head_f,fill=col); d.text((155,y),label,font=small_f,fill=MUTED); y+=38
    # right output excerpt
    out_card=(630,190,W-70,930)
    d.rounded_rectangle(out_card, radius=18, fill=BG, outline=GRID, width=2)
    d.text((665,225),'Output tail excerpt',font=head_f,fill=FG)
    tail=str(r.tail_500 or '')
    gt=str(r.gt_norm)
    # make GT visible in excerpt without claiming semantics
    tail=tail.replace(gt, f'⟦{gt}⟧')
    tail=tail.replace('Final Answer', '\nFinal Answer')
    tail=tail.replace('answer', 'answer')
    y=275
    y=draw_wrapped(d, tail[-900:], (665,y), mono_f, FG, width_chars=92, line_gap=5)
    # footer interpretation
    d.rounded_rectangle((70,960,W-70,1030), radius=14, fill='#EFE7D8', outline=GRID, width=1)
    d.text((100,980), f"Label: {r.pathology_label}. The GT appears in answer-like regions, but the evaluator-visible final commitment drifts to {r.final_pred_norm}.", font=small_f, fill=FG)
    path=OUT/f'08_label_birth_idx{idx:02d}_{run}.png'
    img.save(path)
    return path


def write_index(paths):
    lines=['# KVFidelity Casey v4', '', 'Casey Reas intervention: main atlas + appendix, lifecycle strips, enriched case strips, label-birth panel.', '']
    for p in sorted(OUT.glob('*.png')):
        lines.append(f'- [{p.name}]({p.name})')
    lines += ['', 'Source:', f'- `{SRC}`']
    (OUT/'README.md').write_text('\n'.join(lines), encoding='utf-8')


def add_title(slide, title, subtitle=None):
    tx=slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.7))
    p=tx.text_frame.paragraphs[0]
    p.text=title; p.font.name='Inter'; p.font.size=Pt(28); p.font.bold=True; p.font.color.rgb=RGBColor(21,21,21)
    if subtitle:
        st=slide.shapes.add_textbox(Inches(0.52), Inches(0.88), Inches(12.1), Inches(0.42))
        q=st.text_frame.paragraphs[0]
        q.text=subtitle; q.font.name='Inter'; q.font.size=Pt(11); q.font.color.rgb=RGBColor(111,106,96)


def add_text(slide, lines, left=0.75, top=1.7, width=11.8, height=4.8, size=22):
    box=slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf=box.text_frame
    for i,line in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.text=line; p.font.name='Inter'; p.font.size=Pt(size if i<2 else max(16,size-5)); p.font.color.rgb=RGBColor(21,21,21) if i<2 else RGBColor(111,106,96)
        if i==0: p.font.bold=True


def add_image_centered(slide, path, top=1.25, max_width=12.35, max_height=5.85):
    with Image.open(path) as im:
        w_px, h_px = im.size
    aspect = w_px / h_px
    width = max_width; height = width / aspect
    if height > max_height:
        height = max_height; width = height * aspect
    left = (13.333 - width) / 2
    slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width), height=Inches(height))


def create_pptx(df):
    if Presentation is None:
        return None
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5); blank=prs.slide_layouts[6]
    def bg(slide):
        fill=slide.background.fill; fill.solid(); fill.fore_color.rgb=RGBColor(247,243,234)
    # cover
    s=prs.slides.add_slide(blank); bg(s)
    add_title(s,'KVFidelity Trace Atlas v4','Casey Reas intervention · CASK AIME24 n=30 · Qwen3-8B')
    add_text(s,[
        'A correct answer is not a score. It is a temporal event.',
        'It emerges, persists, drifts, closes, or disappears.',
        'CASK is the experiment. KVFidelity is the lens.',
        'Main deck first. Defensive evidence moves to appendix.'
    ], size=25)
    # main deck
    main_slides=[
        ('Topology matrix','The answer becomes a field: distribution, recurrence, absence, drift.','01_topology_matrix.png',0.95,6.15),
        ('Answer lifecycle','Node labels are counts out of 30; orange bands mark loss between lifecycle stages.','02_answer_lifecycle_strips.png',1.15,5.65),
        ('FullKV 2048 → 4096','More generation transforms states: it repairs some samples and destabilizes others.','04_sankey_fullkv_2048_to_4096.png',1.1,5.7),
        ('Case idx 7','Compression island: several compressed runs close cleanly while FullKV 4096 drifts.','case_strip_idx_07.png',1.1,5.7),
        ('Case idx 9','The answer appears, persists, and still fails to close.','case_strip_idx_09.png',1.1,5.7),
        ('Case idx 11','More generation repairs FullKV closure; compression does not recover it here.','case_strip_idx_11.png',1.1,5.7),
        ('How the label is born','Make the measurement process visible: output excerpt → signals → pathology label.','08_label_birth_idx09_fullkv_4096.png',1.15,5.8),
    ]
    for title,sub,img,top,h in main_slides:
        s=prs.slides.add_slide(blank); bg(s); add_title(s,title,sub); add_image_centered(s,OUT/img,top,max_width=12.35,max_height=h)
    # thesis closing
    s=prs.slides.add_slide(blank); bg(s)
    add_title(s,'Thesis')
    add_text(s,[
        'KVFidelity becomes strong when the correct answer stops being a score.',
        'It becomes an event with a temporal shape.',
        'The useful object is not a leaderboard. It is a score/partiture of answer trajectories.',
        'Public claim remains deferred until manual label audit.'
    ], size=25)
    # appendix divider
    s=prs.slides.add_slide(blank); bg(s); add_title(s,'Appendix','Counts, caveats, alternate transitions, and full case strips.')
    add_text(s,['Appendix: evidence and defensive material','The main deck shows behavior. The appendix keeps the accounting.'], size=25)
    appendix=[
        ('Appendix · D/R/C bars','Conventional count view retained for auditability.','A01_drc_bars.png',1.1,5.7),
        ('Appendix · pathology mix','Raw accuracy collapses these distinct failure modes.','A02_pathology_stack.png',1.1,5.7),
        ('Appendix · TriAttention transition','Transition structure for the TriAttention baseline.','A03_sankey_tri_b256_2048_to_4096.png',1.1,5.7),
        ('Appendix · CASK transition','CASK preserves idx 7, but does not recover FullKV gains here.','A04_sankey_cask_b256_2048_to_4096.png',1.1,5.7),
        ('Appendix · candidate churn','Longer generation opens space for candidate drift.','A05_candidate_churn_scatter.png',1.1,5.7),
        ('Appendix · idx 0','A FullKV 4096 gain: discovery stabilizes into closure.','case_strip_idx_00.png',1.1,5.7),
        ('Appendix · idx 12','A drift-heavy case: churn without clean commitment.','case_strip_idx_12.png',1.1,5.7),
        ('Appendix · idx 24','A FullKV 4096 closure not recovered by compression.','case_strip_idx_24.png',1.1,5.7),
        ('Appendix · idx 26','A FullKV 4096 closure not recovered by compression.','case_strip_idx_26.png',1.1,5.7),
    ]
    for title,sub,img,top,h in appendix:
        s=prs.slides.add_slide(blank); bg(s); add_title(s,title,sub); add_image_centered(s,OUT/img,top,max_width=12.35,max_height=h)
    out=OUT/'kvfidelity_casey_atlas_v4.pptx'
    prs.save(out)
    return out


def make_contact_sheet():
    imgs=[]
    for p in sorted(OUT.glob('*.png')):
        if p.name == 'contact_sheet.png':
            continue
        im=Image.open(p).convert('RGB'); im.thumbnail((400,300))
        canvas=Image.new('RGB',(400,330),'white')
        canvas.paste(im,((400-im.width)//2,0))
        d=ImageDraw.Draw(canvas); d.text((8,307),p.name,fill=(0,0,0))
        imgs.append(canvas)
    cols=2; rows=(len(imgs)+1)//2
    sheet=Image.new('RGB',(cols*400,rows*330),(240,240,240))
    for i,im in enumerate(imgs): sheet.paste(im,((i%cols)*400,(i//cols)*330))
    sheet.save(OUT/'contact_sheet.png')


def main():
    df=load()
    paths=[]
    paths.append(topology(df))
    paths.append(drc_cycle(df))
    paths.append(drc_summary_appendix(df))
    paths.append(pathology_stacked(df))
    paths += transitions(df)
    paths.append(scatter(df))
    paths += make_case_strips(df)
    paths.append(label_birth_panel(df, idx=9, run='fullkv_4096'))
    deck=create_pptx(df)
    make_contact_sheet()
    write_index(paths)
    print(json.dumps({'out':str(OUT),'pngs':len(list(OUT.glob('*.png'))),'deck':str(deck) if deck else None}, indent=2))

if __name__=='__main__':
    main()
