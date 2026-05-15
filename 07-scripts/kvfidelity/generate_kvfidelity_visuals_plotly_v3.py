#!/usr/bin/env python3
import json, os, math
from collections import Counter, defaultdict
from pathlib import Path

import pandas as pd
import plotly.graph_objects as go
from plotly.subplots import make_subplots
import plotly.express as px

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
except Exception:
    Presentation = None

SRC = Path('/home/aya/implante/tmp/kvfidelity-aime24-n30-traces-2026-05-13.jsonl')
OUT = Path('/home/aya/implante/tmp/kvfidelity-plotly-v3-2026-05-13')
OUT.mkdir(parents=True, exist_ok=True)

RUN_ORDER = [
    'fullkv_2048','tri_b256_2048','cask_b256_2048',
    'fullkv_4096','tri_b256_4096','cask_b256_4096',
    'tri_b384_4096','cask_b384_4096','tri_b512_4096','cask_b512_4096'
]
IMPORTANT = [0,7,9,11,12,24,26]

LABELS = {
    'clean_correct': {'sym':'✓','color':'#007A4D','name':'closed clean'},
    'correct_unboxed': {'sym':'u','color':'#69B99A','name':'closed weak'},
    'answer_marker_drift': {'sym':'D','color':'#D65A3B','name':'answer drift'},
    'latent_answer_marker_not_closed': {'sym':'A','color':'#C98713','name':'latent after marker'},
    'latent_final_zone_not_closed': {'sym':'F','color':'#E3B341','name':'latent final zone'},
    'latent_intermediate_only': {'sym':'~','color':'#7655B7','name':'intermediate only'},
    'evaluator_missed_boxed_gt': {'sym':'B','color':'#368CBF','name':'boxed missed'},
    'no_gt_seen': {'sym':'·','color':'#D6D1C7','name':'not discovered'},
}
LABEL_ORDER = list(LABELS.keys())
LABEL_VALUE = {k:i for i,k in enumerate(LABEL_ORDER)}
COLORSCALE = []
n=len(LABEL_ORDER)
for k,i in LABEL_VALUE.items():
    c=LABELS[k]['color']
    lo=i/(n-1) if n>1 else 0
    COLORSCALE.append([lo,c])
    COLORSCALE.append([lo,c])

BG = '#FFFDF7'
PAPER = '#F7F3EA'
FG = '#151515'
MUTED = '#6F6A60'
GRID = '#D8D0C2'
ACCENT = '#007A4D'
ORANGE = '#D65A3B'
PURPLE = '#7655B7'
BLUE = '#368CBF'
YELLOW = '#C98713'


def load():
    rows=[]
    for line in SRC.read_text().splitlines():
        if line.strip(): rows.append(json.loads(line))
    df=pd.DataFrame(rows)
    df['run_id'] = pd.Categorical(df['run_id'], categories=RUN_ORDER, ordered=True)
    df['closed'] = df['eval_correct'].astype(bool)
    df['discovered'] = df['gt_anywhere'].astype(bool)
    df['retained'] = (df['gt_in_final_zone'].astype(bool) | df['gt_after_answer_marker'].astype(bool) | df['gt_boxed'].astype(bool))
    df['run_short'] = df['run_id'].astype(str).str.replace('_','<br>')
    return df.sort_values(['idx','run_id'])


def style(fig, title=None, w=1600, h=900):
    fig.update_layout(
        template='plotly_white',
        paper_bgcolor=PAPER,
        plot_bgcolor=BG,
        font=dict(family='Inter, Arial, sans-serif', color=FG, size=18),
        title=dict(text=title or '', x=0.02, y=0.98, xanchor='left', yanchor='top', font=dict(size=34, color=FG)),
        margin=dict(l=80,r=60,t=100,b=80),
        width=w,
        height=h,
    )
    fig.update_xaxes(gridcolor=GRID, zerolinecolor=GRID, tickfont=dict(color=MUTED, size=14))
    fig.update_yaxes(gridcolor=GRID, zerolinecolor=GRID, tickfont=dict(color=MUTED, size=14))
    return fig


def save(fig, name, w=1600, h=900):
    html=OUT/f'{name}.html'
    png=OUT/f'{name}.png'
    fig.write_html(str(html), include_plotlyjs='cdn')
    fig.write_image(str(png), width=w, height=h, scale=2)
    return png


def topology(df):
    mat=[]; text=[]; hover=[]
    for idx in range(30):
        row=[]; trow=[]; hrow=[]
        for run in RUN_ORDER:
            r=df[(df.idx==idx)&(df.run_id.astype(str)==run)].iloc[0]
            lab=r.pathology_label
            row.append(LABEL_VALUE[lab]); trow.append(LABELS[lab]['sym'])
            hrow.append(f"idx {idx}<br>GT {r.gt_norm}<br>{run}<br>{LABELS[lab]['name']}<br>pred {r.final_pred_norm}<br>first {r.first_gt_frac}<br>last {r.last_gt_frac}<br>flips {r.candidate_flips}")
        mat.append(row); text.append(trow); hover.append(hrow)
    fig=go.Figure(go.Heatmap(
        z=mat, x=[r.replace('_','<br>') for r in RUN_ORDER], y=list(range(30)),
        colorscale=COLORSCALE, showscale=False, text=text, texttemplate='%{text}',
        textfont=dict(size=22, color='#151515'), hovertext=hover, hoverinfo='text',
        xgap=3, ygap=2
    ))
    fig=style(fig, 'Topology matrix: the life of the answer across runs', 1800, 1450)
    fig.update_layout(margin=dict(l=95,r=360,t=120,b=150))
    fig.update_yaxes(autorange='reversed', title='AIME24 problem index')
    fig.update_xaxes(title='runtime condition', tickfont=dict(size=13, color=MUTED), title_standoff=28)
    # Legend in the right margin to avoid collision with x-axis labels.
    fig.add_annotation(
        text='Legend', x=1.025, y=0.98, xref='paper', yref='paper', showarrow=False,
        xanchor='left', align='left', font=dict(size=17,color=FG)
    )
    for i,(k,v) in enumerate(LABELS.items()):
        fig.add_annotation(
            text=f"<span style='color:{v['color']}; font-size:22px'>{v['sym']}</span> {v['name']}",
            x=1.025, y=0.925 - i*0.055, xref='paper', yref='paper', showarrow=False,
            xanchor='left', align='left', font=dict(size=15,color=MUTED)
        )
    return save(fig,'01_topology_matrix',1800,1450)


def drc_summary(df):
    runs=[]; discovery=[]; retention=[]; closure=[]
    for run in RUN_ORDER:
        sub=df[df.run_id.astype(str)==run]
        runs.append(run.replace('_','<br>'))
        discovery.append(sub.discovered.mean()*100)
        retention.append(sub.retained.mean()*100)
        closure.append(sub.closed.mean()*100)
    fig=go.Figure()
    fig.add_bar(name='Discovery: GT appears', x=runs, y=discovery, marker_color=PURPLE)
    fig.add_bar(name='Retention: GT near final/marker', x=runs, y=retention, marker_color=YELLOW)
    fig.add_bar(name='Closure: eval correct', x=runs, y=closure, marker_color=ACCENT)
    fig=style(fig,'Discovery / Retention / Closure are different instruments',1800,950)
    fig.update_layout(barmode='group', legend=dict(orientation='h', y=1.05, x=0.02), yaxis_ticksuffix='%')
    fig.update_yaxes(range=[0,45], title='share of 30 problems')
    fig.add_annotation(text='FullKV 4096 discovers 11/30 but closes only 4/30. Compression mostly reduces discovery; when it discovers, closure can still fail.', x=0.02, y=-0.18, xref='paper', yref='paper', showarrow=False, align='left', font=dict(size=22,color=FG))
    return save(fig,'02_discovery_retention_closure',1800,950)


def pathology_stacked(df):
    counts=[]
    for run in RUN_ORDER:
        sub=df[df.run_id.astype(str)==run]
        c=Counter(sub.pathology_label)
        for lab in LABEL_ORDER:
            counts.append({'run':run.replace('_','<br>'),'label':LABELS[lab]['name'],'count':c.get(lab,0),'color':LABELS[lab]['color']})
    fig=go.Figure()
    for lab in LABEL_ORDER:
        name=LABELS[lab]['name']; color=LABELS[lab]['color']
        ys=[]
        for run in RUN_ORDER:
            sub=df[df.run_id.astype(str)==run]
            ys.append((sub.pathology_label==lab).sum())
        fig.add_bar(name=name, x=[r.replace('_','<br>') for r in RUN_ORDER], y=ys, marker_color=color)
    fig=style(fig,'Pathology mix: accuracy collapses several failure modes',1800,950)
    fig.update_layout(
        barmode='stack',
        margin=dict(l=90,r=50,t=135,b=175),
        legend=dict(
            orientation='h',
            y=-0.22,
            x=0.0,
            xanchor='left',
            yanchor='top',
            font=dict(size=14),
            traceorder='normal',
            itemwidth=120,
        )
    )
    fig.update_yaxes(title='count of 30 samples')
    return save(fig,'03_pathology_stack',1800,950)


def hex_to_rgba(hex_color, alpha=0.64):
    h=hex_color.lstrip('#')
    r=int(h[0:2],16); g=int(h[2:4],16); b=int(h[4:6],16)
    return f'rgba({r},{g},{b},{alpha})'

def sankey_for(df, a, b, name):
    labels=['not_discovered','discovered_not_retained','closure_failure','closed']
    colors={'not_discovered':'#D6D1C7','discovered_not_retained':PURPLE,'closure_failure':ORANGE,'closed':ACCENT}
    def cls(pathology_class): return pathology_class
    left=[f'2048<br>{l}' for l in labels]
    right=[f'4096<br>{l}' for l in labels]
    node_labels=left+right
    node_colors=[colors[l] for l in labels]+[colors[l] for l in labels]
    trans=Counter()
    for idx in range(30):
        ra=df[(df.idx==idx)&(df.run_id.astype(str)==a)].iloc[0]
        rb=df[(df.idx==idx)&(df.run_id.astype(str)==b)].iloc[0]
        trans[(ra.pathology_class, rb.pathology_class)] += 1
    source=[]; target=[]; value=[]; link_color=[]
    for (x,y),v in trans.items():
        source.append(labels.index(x)); target.append(len(labels)+labels.index(y)); value.append(v); link_color.append(hex_to_rgba(colors[y], 0.62))
    fig=go.Figure(go.Sankey(
        arrangement='fixed',
        node=dict(label=node_labels, color=node_colors, pad=18, thickness=24, line=dict(color='#FFFDF7',width=1)),
        link=dict(source=source,target=target,value=value,color=link_color)
    ))
    display_titles = {
        '04_sankey_fullkv_2048_to_4096': 'FullKV: 2048 → 4096 transition',
        '05_sankey_tri_b256_2048_to_4096': 'TriAttention b256: 2048 → 4096 transition',
        '06_sankey_cask_b256_2048_to_4096': 'CASK b256: 2048 → 4096 transition',
    }
    fig=style(fig, display_titles.get(name, name), 1500, 900)
    return save(fig, name,1500,900)


def transitions(df):
    return [
        sankey_for(df,'fullkv_2048','fullkv_4096','04_sankey_fullkv_2048_to_4096'),
        sankey_for(df,'tri_b256_2048','tri_b256_4096','05_sankey_tri_b256_2048_to_4096'),
        sankey_for(df,'cask_b256_2048','cask_b256_4096','06_sankey_cask_b256_2048_to_4096'),
    ]


def scatter(df):
    fig=go.Figure()
    for lab in LABEL_ORDER:
        sub=df[df.pathology_label==lab]
        if sub.empty: continue
        fig.add_trace(go.Scatter(
            x=sub.output_chars, y=sub.candidate_flips, mode='markers', name=LABELS[lab]['name'],
            text=[f"{r.run_id}<br>idx {r.idx}<br>gt {r.gt_norm}<br>{LABELS[lab]['name']}" for r in sub.itertuples()],
            hoverinfo='text+x+y', marker=dict(size=12, color=LABELS[lab]['color'], line=dict(width=1,color='#000'))
        ))
    fig=style(fig,'Candidate churn vs output length: more generation can mean more drift',1600,950)
    fig.update_xaxes(title='output characters')
    fig.update_yaxes(title='candidate flips near answer markers')
    return save(fig,'07_candidate_churn_scatter',1600,950)


def case_strip(df, idx):
    sub=df[df.idx==idx].copy()
    fig=go.Figure()
    ylabels=[r.replace('_',' ') for r in RUN_ORDER]
    ymap={run:i for i,run in enumerate(RUN_ORDER)}
    # baseline
    for run in RUN_ORDER:
        y=ymap[run]
        fig.add_trace(go.Scatter(x=[0,1],y=[y,y],mode='lines',line=dict(color=GRID,width=6),showlegend=False,hoverinfo='skip'))
    for _,r in sub.iterrows():
        y=ymap[str(r.run_id)]
        lab=r.pathology_label
        color=LABELS[lab]['color']
        if pd.notna(r.first_gt_frac):
            fig.add_trace(go.Scatter(x=[r.first_gt_frac],y=[y],mode='markers',name='first GT',marker=dict(symbol='triangle-up',size=16,color=color),showlegend=False,hovertext=f"{r.run_id}<br>first GT {r.first_gt_frac}",hoverinfo='text'))
        if pd.notna(r.last_gt_frac):
            fig.add_trace(go.Scatter(x=[r.last_gt_frac],y=[y],mode='markers',name='last GT',marker=dict(symbol='diamond',size=16,color=color),showlegend=False,hovertext=f"{r.run_id}<br>last GT {r.last_gt_frac}",hoverinfo='text'))
        # closure marker at end
        fig.add_trace(go.Scatter(x=[1.03],y=[y],mode='markers+text',text=[LABELS[lab]['sym']],textposition='middle center',marker=dict(symbol='square',size=28,color=color),textfont=dict(size=18,color='#151515'),showlegend=False,hovertext=f"{r.run_id}<br>{LABELS[lab]['name']}<br>pred {r.final_pred_norm}<br>flips {r.candidate_flips}",hoverinfo='text'))
    gt=sub.iloc[0].gt_norm
    fig=style(fig,f'Case strip idx {idx}: GT {gt}',1700,900)
    fig.update_xaxes(range=[-0.02,1.08], title='normalized position in output', tickformat='.0%')
    fig.update_yaxes(tickmode='array', tickvals=list(range(len(RUN_ORDER))), ticktext=ylabels, autorange='reversed')
    fig.add_annotation(text='▲ first GT · ◆ last GT · square = closure/pathology', x=0.02, y=-0.15, xref='paper', yref='paper', showarrow=False, align='left', font=dict(size=22,color=MUTED))
    return save(fig,f'case_strip_idx_{idx:02d}',1700,900)


def make_case_strips(df):
    return [case_strip(df,idx) for idx in IMPORTANT]


def write_index(paths):
    lines=['# KVFidelity Plotly Atlas v2', '', 'Plotly/Kaleido light-theme visuals generated from trace JSONL.', '']
    for p in sorted(OUT.glob('*.png')):
        lines.append(f'- [{p.name}]({p.name})')
    lines += ['', 'Source:', f'- `{SRC}`', f'- `{OUT}/kvfidelity_plotly_atlas_v2.pptx`']
    (OUT/'README.md').write_text('\n'.join(lines), encoding='utf-8')


def add_title(slide, title, subtitle=None):
    tx=slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.7))
    p=tx.text_frame.paragraphs[0]
    p.text=title
    p.font.name='Inter'
    p.font.size=Pt(28)
    p.font.bold=True
    p.font.color.rgb=RGBColor(21,21,21)
    if subtitle:
        st=slide.shapes.add_textbox(Inches(0.52), Inches(0.88), Inches(12.1), Inches(0.35))
        q=st.text_frame.paragraphs[0]
        q.text=subtitle
        q.font.name='Inter'
        q.font.size=Pt(11)
        q.font.color.rgb=RGBColor(111,106,96)


def add_image_centered(slide, path, top=1.25, max_width=12.35, max_height=5.85):
    from PIL import Image
    with Image.open(path) as im:
        w_px, h_px = im.size
    aspect = w_px / h_px
    width = max_width
    height = width / aspect
    if height > max_height:
        height = max_height
        width = height * aspect
    left = (13.333 - width) / 2
    slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width), height=Inches(height))


def create_pptx(df):
    if Presentation is None:
        return None
    prs=Presentation()
    prs.slide_width=Inches(13.333)
    prs.slide_height=Inches(7.5)
    blank=prs.slide_layouts[6]
    def bg(slide):
        fill=slide.background.fill
        fill.solid(); fill.fore_color.rgb=RGBColor(247,243,234)
    # Title
    s=prs.slides.add_slide(blank); bg(s)
    add_title(s,'KVFidelity Trace Atlas v2','CASK AIME24 n=30 · Qwen3-8B · FullKV / TriAttention / CASK')
    box=s.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.8), Inches(3.8))
    tf=box.text_frame
    for i,line in enumerate([
        'CASK is the experiment. KVFidelity is the lens.',
        'The object is not a leaderboard: it is the life of the answer.',
        'Discovery: does the answer emerge?',
        'Retention: does it remain accessible?',
        'Closure: does it become a clean final commitment?'
    ]):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.text=line; p.font.name='Inter'; p.font.size=Pt(26 if i<2 else 19); p.font.color.rgb=RGBColor(21,21,21) if i<2 else RGBColor(111,106,96)
        if i==0: p.font.bold=True
    slides=[
        ('Topology matrix','Samples × conditions: the pattern appears by trajectory, not by average.','01_topology_matrix.png',0.95,6.15),
        ('Discovery / Retention / Closure','FullKV 4096 discovers 11/30 and closes 4/30. Compression reduces discovery and stability.','02_discovery_retention_closure.png',1.1,5.7),
        ('Pathology mix','Raw accuracy collapses distinct failure modes.','03_pathology_stack.png',1.1,5.7),
        ('FullKV 2048 → 4096','More tokens improve aggregate accuracy, but can move closed cases into drift.','04_sankey_fullkv_2048_to_4096.png',1.1,5.7),
        ('TriAttention b256 2048 → 4096','Transition structure for the TriAttention baseline.','05_sankey_tri_b256_2048_to_4096.png',1.1,5.7),
        ('CASK b256 2048 → 4096','CASK preserves idx 7, but does not recover FullKV gains.','06_sankey_cask_b256_2048_to_4096.png',1.1,5.7),
        ('Candidate churn','Longer generation opens space for candidate drift.','07_candidate_churn_scatter.png',1.1,5.7),
        ('Case strip: idx 0','A FullKV 4096 gain: discovery stabilizes into closure.','case_strip_idx_00.png',1.1,5.7),
        ('Case strip: idx 7','A stable compression island: CASK/Tri b512 close, FullKV 4096 drifts.','case_strip_idx_07.png',1.1,5.7),
        ('Case strip: idx 9','The answer appears, but closure fails under long generation/compression.','case_strip_idx_09.png',1.1,5.7),
        ('Case strip: idx 11','A closure_failure → closed example in FullKV 4096; compression does not recover it.','case_strip_idx_11.png',1.1,5.7),
        ('Case strip: idx 12','A drift-heavy case: candidate churn without clean commitment.','case_strip_idx_12.png',1.1,5.7),
        ('Case strip: idx 24','A FullKV 4096 closure not recovered by compression.','case_strip_idx_24.png',1.1,5.7),
        ('Case strip: idx 26','A FullKV 4096 closure not recovered by compression.','case_strip_idx_26.png',1.1,5.7),
    ]
    for title,sub,img,top,h in slides:
        s=prs.slides.add_slide(blank); bg(s); add_title(s,title,sub); add_image_centered(s,OUT/img,top,max_width=12.35,max_height=h)
    out=OUT/'kvfidelity_plotly_atlas_v3.pptx'
    prs.save(out)
    return out


def main():
    df=load()
    paths=[]
    for fn in [topology, drc_summary, pathology_stacked, scatter]:
        paths.append(fn(df))
    paths += transitions(df)
    paths += make_case_strips(df)
    deck=create_pptx(df)
    write_index(paths)
    print(json.dumps({'out':str(OUT),'pngs':len(list(OUT.glob('*.png'))),'deck':str(deck) if deck else None}, indent=2))

if __name__=='__main__':
    main()
